from fastapi import FastAPI, Depends, HTTPException, Header
from fastapi.responses import StreamingResponse
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from typing import Optional
from groq import Groq
from supabase import create_client, Client
import os
import base64
from io import BytesIO
from dotenv import load_dotenv

load_dotenv()

app = FastAPI()

allowed_origins = [
    o.strip()
    for o in os.getenv(
        "ALLOWED_ORIGINS",
        "https://zenithgpt.vercel.app,http://localhost:3000",
    ).split(",")
    if o.strip()
]

app.add_middleware(
    CORSMiddleware,
    allow_origins=allowed_origins,
    allow_methods=["GET", "POST", "PATCH", "DELETE"],
    allow_headers=["Authorization", "Content-Type"],
)

groq_client = Groq(api_key=os.getenv("GROQ_API_KEY"))
supabase: Client = create_client(os.getenv("SUPABASE_URL"), os.getenv("SUPABASE_KEY"))

TEXT_MODEL = os.getenv("GROQ_TEXT_MODEL", "llama-3.3-70b-versatile")
VISION_MODEL = os.getenv("GROQ_VISION_MODEL", "meta-llama/llama-4-scout-17b-16e-instruct")


def get_current_user_id(authorization: Optional[str] = Header(None)) -> str:
    if not authorization or not authorization.startswith("Bearer "):
        raise HTTPException(status_code=401, detail="Missing bearer token")
    token = authorization.split(" ", 1)[1]
    try:
        user_response = supabase.auth.get_user(token)
    except Exception:
        raise HTTPException(status_code=401, detail="Invalid token")
    user = getattr(user_response, "user", None)
    if not user or not getattr(user, "id", None):
        raise HTTPException(status_code=401, detail="Invalid token")
    return user.id


class Attachment(BaseModel):
    type: str
    name: str
    data: str


class ChatRequest(BaseModel):
    message: str
    mode: str = "general"
    history: list = []
    attachments: list[Attachment] = []


class ChatResponse(BaseModel):
    reply: str
    error: Optional[str] = None


class SaveChatRequest(BaseModel):
    conversation_id: Optional[str] = None
    title: str
    messages: list
    mode: str = "general"


class UpdateChatRequest(BaseModel):
    title: Optional[str] = None
    pinned: Optional[bool] = None
    archived: Optional[bool] = None


ATTACHMENT_GUIDANCE = (
    "\n\nWhen the user uploads a file, its extracted contents appear in their message "
    "between '--- Attached file: <name> ---' and '--- end <name> ---' markers. Treat "
    "that text as the file's actual contents — read it directly and answer based on it. "
    "Never say 'I can't access files' or 'I can't open attachments' when these markers "
    "are present; the content is already in front of you. If extraction returned an "
    "error placeholder (starts with '[Error' or '[Could not'), say so honestly and ask "
    "the user to re-upload or describe what's in the file."
)

SYSTEM_PROMPTS = {
    "general": (
        "You are ZenithGPT, a thoughtful and direct AI assistant.\n\n"
        "You can read images, PDFs, Word documents (.docx), and PowerPoint files (.pptx) "
        "when users upload them — the contents are extracted and provided to you. "
        "If a user asks whether you can read these formats, the answer is yes.\n\n"
        "Match the depth of your reply to the question — short questions get short answers, "
        "complex ones get structured ones. Use markdown (headers, lists, code blocks) when it "
        "aids clarity, not as decoration. Skip filler phrases like 'Great question!' or 'I hope this helps!'.\n\n"
        "When a request is ambiguous, ask one focused clarifying question instead of guessing. "
        "When you don't know something, say so plainly. When you're confident, be confident — "
        "don't hedge for politeness. If a user pushes back with a correction, weigh it on its merits "
        "rather than capitulating reflexively.\n\n"
        "Format code in fenced blocks with the language tag. Use concrete examples for math, "
        "science, or abstract reasoning. End answers when they're complete — no summary paragraph "
        "restating what you just said."
    ),
    "health": (
        "You are ZenithGPT in health mode — a knowledgeable medical assistant.\n\n"
        "Provide accurate, evidence-based information. Explain the mechanism (why something works), "
        "not just the recommendation. Use plain language and define medical terms in parentheses on "
        "first use. Be specific about dosages, timeframes, and warning signs when relevant.\n\n"
        "For serious symptoms or medication questions, add: 'I'm not a substitute for a doctor — "
        "please consult one for anything serious.' Don't pile on disclaimers for general wellness "
        "questions; one reminder is enough.\n\n"
        "If a symptom could be urgent (chest pain, sudden severe headache, difficulty breathing, "
        "stroke signs, severe allergic reaction, etc.), tell them to seek emergency care immediately, "
        "before anything else."
    ),
    "code": (
        "You are ZenithGPT in code mode — an expert programming assistant working alongside the user "
        "like a senior engineer.\n\n"
        "Write working, idiomatic code with clear names. Default to no comments — add one only when "
        "the *why* is non-obvious. Use the user's stated language and stack; if unspecified, ask "
        "before assuming.\n\n"
        "When debugging, find the root cause rather than patching symptoms. When designing, name the "
        "main trade-off so the user can decide. Use fenced code blocks with language tags. Keep "
        "explanations tight — show the code, briefly say what changed and why.\n\n"
        "If a request is ambiguous (which framework? sync or async? what version?), ask one quick "
        "question rather than producing five variations."
    ),
    "writing": (
        "You are ZenithGPT in writing mode — a sharp writing partner.\n\n"
        "Match the user's voice and intent. Cut filler, weak verbs, and unnecessary adverbs. Vary "
        "sentence length for rhythm. For business writing, be direct; for creative, be vivid; for "
        "academic, be precise.\n\n"
        "When the user shares a draft, name specifically what works and what doesn't, then offer "
        "concrete revisions — not vague suggestions like 'consider rewording'. When they ask for new "
        "writing, deliver clean prose without preamble. 'Here's a draft:' is fine; 'I hope this "
        "helps!' is not."
    ),
    "data": (
        "You are ZenithGPT in data mode — an experienced data scientist and analyst.\n\n"
        "Recommend the simplest method that answers the question. Don't reach for machine learning "
        "when SQL or a chart will do. State your assumptions explicitly: sample size, distribution, "
        "missing data, sources of bias.\n\n"
        "For code, prefer pandas, polars, SQL, or numpy as appropriate, with small reproducible "
        "examples. When interpreting results, translate the number into plain terms (e.g. "
        "'an effect size of 0.3 means about a third of a standard deviation'). Distinguish "
        "correlation from causation when it matters. Flag datasets that are too small, biased, "
        "or missing key context."
    ),
}


MAX_EXTRACTED_CHARS = 50000


def _decode_data_url(data_url: str) -> bytes:
    if "," in data_url:
        data_url = data_url.split(",", 1)[1]
    return base64.b64decode(data_url)


def extract_text_from_file(name: str, data_url: str) -> str:
    ext = name.rsplit(".", 1)[-1].lower() if "." in name else ""
    try:
        binary = _decode_data_url(data_url)
    except Exception as e:
        return f"[Could not decode {name}: {e}]"

    try:
        if ext == "pdf":
            from pypdf import PdfReader

            reader = PdfReader(BytesIO(binary))
            pages = [p.extract_text() or "" for p in reader.pages]
            return "\n\n".join(f"Page {i+1}:\n{t}" for i, t in enumerate(pages) if t.strip())
        if ext == "docx":
            from docx import Document

            doc = Document(BytesIO(binary))
            return "\n".join(p.text for p in doc.paragraphs if p.text)
        if ext == "pptx":
            from pptx import Presentation

            prs = Presentation(BytesIO(binary))
            slides = []
            for i, slide in enumerate(prs.slides, 1):
                parts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        parts.append(shape.text)
                if parts:
                    slides.append(f"Slide {i}:\n" + "\n".join(parts))
            return "\n\n".join(slides)
        if ext == "xlsx":
            from openpyxl import load_workbook

            wb = load_workbook(BytesIO(binary), data_only=True, read_only=True)
            sheets = []
            for ws in wb.worksheets:
                rows = []
                for row in ws.iter_rows(values_only=True):
                    rows.append("\t".join("" if c is None else str(c) for c in row))
                sheets.append(f"Sheet: {ws.title}\n" + "\n".join(rows))
            return "\n\n".join(sheets)
    except Exception as e:
        return f"[Error extracting {name}: {e}]"

    return f"[File {name} attached but its format ({ext or 'unknown'}) is not supported for text extraction.]"


def build_messages(request: ChatRequest):
    base = SYSTEM_PROMPTS.get(request.mode, SYSTEM_PROMPTS["general"])
    has_attachments = bool(request.attachments)
    system_message = base + (ATTACHMENT_GUIDANCE if has_attachments else "")
    messages = [{"role": "system", "content": system_message}]
    messages.extend(request.history)

    images = [a for a in request.attachments if a.type == "image"]

    text_part = request.message
    for a in request.attachments:
        if a.type == "image":
            continue
        if a.type == "file":
            content = extract_text_from_file(a.name, a.data)
        else:
            content = a.data
        if len(content) > MAX_EXTRACTED_CHARS:
            content = content[:MAX_EXTRACTED_CHARS] + "\n[...truncated]"
        text_part += f"\n\n--- Attached file: {a.name} ---\n{content}\n--- end {a.name} ---"

    if images:
        content = [{"type": "text", "text": text_part}]
        for img in images:
            url = img.data if img.data.startswith("data:") else f"data:image/jpeg;base64,{img.data}"
            content.append({"type": "image_url", "image_url": {"url": url}})
        messages.append({"role": "user", "content": content})
    else:
        messages.append({"role": "user", "content": text_part})

    return messages


def pick_model(request: ChatRequest) -> str:
    has_image = any(a.type == "image" for a in request.attachments)
    return VISION_MODEL if has_image else TEXT_MODEL


@app.get("/")
def read_root():
    return {"message": "ZenithGPT Backend API"}


@app.post("/chat", response_model=ChatResponse)
async def chat(request: ChatRequest, user_id: str = Depends(get_current_user_id)):
    try:
        response = groq_client.chat.completions.create(
            model=pick_model(request),
            messages=build_messages(request),
            max_tokens=4096,
            temperature=0.6,
            top_p=0.9,
        )
        reply = response.choices[0].message.content
        return ChatResponse(reply=reply)
    except Exception as e:
        return ChatResponse(reply="", error=str(e))


@app.post("/chat-stream")
async def chat_stream(request: ChatRequest, user_id: str = Depends(get_current_user_id)):
    messages = build_messages(request)
    model = pick_model(request)

    def generate():
        try:
            stream = groq_client.chat.completions.create(
                model=model,
                messages=messages,
                max_tokens=4096,
                temperature=0.6,
                top_p=0.9,
                stream=True,
            )
            for chunk in stream:
                content = chunk.choices[0].delta.content
                if content:
                    yield content
        except Exception as e:
            yield f"\n\n[Error: {str(e)}]"

    return StreamingResponse(generate(), media_type="text/plain")


@app.post("/save-chat")
async def save_chat(
    request: SaveChatRequest, user_id: str = Depends(get_current_user_id)
):
    try:
        if request.conversation_id:
            existing = (
                supabase.table("conversations")
                .select("user_id")
                .eq("id", request.conversation_id)
                .single()
                .execute()
            )
            if not existing.data or existing.data["user_id"] != user_id:
                raise HTTPException(status_code=403, detail="Not your conversation")
            result = (
                supabase.table("conversations")
                .update({
                    "messages": request.messages,
                    "title": request.title,
                    "mode": request.mode,
                })
                .eq("id", request.conversation_id)
                .eq("user_id", user_id)
                .execute()
            )
        else:
            result = (
                supabase.table("conversations")
                .insert({
                    "title": request.title,
                    "messages": request.messages,
                    "mode": request.mode,
                    "user_id": user_id,
                })
                .execute()
            )
        return {"success": True, "data": result.data}
    except HTTPException:
        raise
    except Exception as e:
        return {"success": False, "error": str(e)}


@app.get("/conversations")
async def get_conversations(user_id: str = Depends(get_current_user_id)):
    try:
        result = (
            supabase.table("conversations")
            .select("*")
            .eq("user_id", user_id)
            .order("updated_at", desc=True)
            .execute()
        )
        rows = [r for r in (result.data or []) if not r.get("archived")]
        return {"conversations": rows}
    except Exception as e:
        return {"conversations": [], "error": str(e)}


@app.patch("/conversations/{conversation_id}")
async def update_conversation(
    conversation_id: str,
    request: UpdateChatRequest,
    user_id: str = Depends(get_current_user_id),
):
    try:
        updates = {}
        if request.title is not None:
            updates["title"] = request.title
        if request.pinned is not None:
            updates["pinned"] = request.pinned
        if request.archived is not None:
            updates["archived"] = request.archived
        result = (
            supabase.table("conversations")
            .update(updates)
            .eq("id", conversation_id)
            .eq("user_id", user_id)
            .execute()
        )
        if not result.data:
            raise HTTPException(status_code=404, detail="Not found")
        return {"success": True}
    except HTTPException:
        raise
    except Exception as e:
        return {"success": False, "error": str(e)}


@app.delete("/conversations/{conversation_id}")
async def delete_conversation(
    conversation_id: str, user_id: str = Depends(get_current_user_id)
):
    try:
        result = (
            supabase.table("conversations")
            .delete()
            .eq("id", conversation_id)
            .eq("user_id", user_id)
            .execute()
        )
        if not result.data:
            raise HTTPException(status_code=404, detail="Not found")
        return {"success": True}
    except HTTPException:
        raise
    except Exception as e:
        return {"success": False, "error": str(e)}


@app.delete("/delete-account")
async def delete_account(user_id: str = Depends(get_current_user_id)):
    try:
        supabase.table("conversations").delete().eq("user_id", user_id).execute()
        try:
            supabase.auth.admin.delete_user(user_id)
        except Exception:
            pass
        return {"success": True}
    except Exception as e:
        return {"success": False, "error": str(e)}


if __name__ == "__main__":
    import uvicorn

    uvicorn.run(app, host="0.0.0.0", port=8000)
